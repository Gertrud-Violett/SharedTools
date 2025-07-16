#2025 Makkiblog.com MIT License
#

import os
import tkinter as tk
from tkinter import filedialog, messagebox
from pptx import Presentation
from pptx.util import Inches
from markdown import markdown


templatedirectory = "C:\mnt\data\ppt_template.pptx"

def convert_marp_to_powerpoint(markdown_file, template_file):
    try:
        # Get the directory and filename of the input file
        input_dir = os.path.dirname(markdown_file)
        input_filename = os.path.basename(markdown_file)
        output_filename = os.path.splitext(input_filename)[0] + ".pptx"
        output_file = os.path.join(input_dir, output_filename)
        
        # Read the Markdown file
        with open(markdown_file, 'r', encoding='utf-8') as file:
            markdown_text = file.read()

        # Convert Markdown to HTML
        html_text = markdown(markdown_text)

        # Load the PowerPoint template
        prs = Presentation(template_file)

        # Split the HTML into slides based on <hr> tags
        slides = html_text.split('<hr />')

        # Iterate over each slide
        for slide_html in slides:
            # Add a new slide using the first layout
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)

            # Add the slide content
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]

            # Extract the slide title and content
            title_start = slide_html.find('<h1>')
            title_end = slide_html.find('</h1>')
            if title_start != -1 and title_end != -1:
                title_text = slide_html[title_start + 4:title_end]
                title_shape.text = title_text
                slide_html = slide_html[title_end + 5:]

            # Add the remaining content to the slide body
            body_shape.text = slide_html.strip()

            # Add images to the slide
            img_start = slide_html.find('<img src="')
            while img_start != -1:
                img_end = slide_html.find('"', img_start + 10)
                img_path = slide_html[img_start + 10:img_end]
                img_path = os.path.join(input_dir, img_path)

                if os.path.exists(img_path):
                    pic = shapes.add_picture(img_path, Inches(1), Inches(1))

                img_start = slide_html.find('<img src="', img_end)

        # Save the PowerPoint presentation
        prs.save(output_file)
        return output_file
    except Exception as e:
        raise Exception(f"Error converting file: {str(e)}")

class MarpToPowerPointConverter:
    def __init__(self, root):
        self.root = root
        self.root.title("Marp to PowerPoint Converter")
        self.root.geometry("500x200")
        
        # Template file path - hardcoded as per requirements
        self.template_file = templatedirectory
        
        # Create UI elements
        self.create_widgets()
    
    def create_widgets(self):
        # Frame for file selection
        file_frame = tk.Frame(self.root, pady=20)
        file_frame.pack(fill=tk.X)
        
        # Label
        tk.Label(file_frame, text="Select Marp Markdown File:").pack(side=tk.LEFT, padx=10)
        
        # Entry for file path
        self.file_entry = tk.Entry(file_frame, width=30)
        self.file_entry.pack(side=tk.LEFT, padx=5, fill=tk.X, expand=True)
        
        # Browse button
        browse_button = tk.Button(file_frame, text="Browse", command=self.browse_file)
        browse_button.pack(side=tk.LEFT, padx=5)
        
        # Convert button
        convert_button = tk.Button(self.root, text="Convert to PowerPoint", 
                                  command=self.convert_file, padx=10, pady=5)
        convert_button.pack(pady=10)
        
        # Status label
        self.status_label = tk.Label(self.root, text="")
        self.status_label.pack(pady=10)
    
    def browse_file(self):
        file_path = filedialog.askopenfilename(
            filetypes=[("Markdown Files", "*.md")]
        )
        if file_path:
            self.file_entry.delete(0, tk.END)
            self.file_entry.insert(0, file_path)
    
    def convert_file(self):
        markdown_file = self.file_entry.get().strip()
        
        if not markdown_file:
            messagebox.showerror("Error", "Please select a Markdown file")
            return
        
        if not os.path.exists(markdown_file):
            messagebox.showerror("Error", "The selected file does not exist")
            return
        
        if not os.path.exists(self.template_file):
            messagebox.showerror("Error", f"Template file '{self.template_file}' not found")
            return
        
        try:
            self.status_label.config(text="Converting...", fg="blue")
            self.root.update()
            
            output_file = convert_marp_to_powerpoint(markdown_file, self.template_file)
            
            self.status_label.config(
                text=f"Successfully converted to {os.path.basename(output_file)}", 
                fg="green"
            )
            
            # Ask if user wants to open the file
            if messagebox.askyesno("Success", 
                                  f"File converted successfully to {output_file}. Open it now?"):
                os.startfile(output_file)
                
        except Exception as e:
            self.status_label.config(text=f"Error: {str(e)}", fg="red")
            messagebox.showerror("Error", str(e))

if __name__ == "__main__":
    root = tk.Tk()
    app = MarpToPowerPointConverter(root)
    root.mainloop()